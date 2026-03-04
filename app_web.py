import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from docx import Document
from io import BytesIO
import os
import re

st.set_page_config(
    page_title="Gerador de Slides da Pauta",
    layout="centered"
)

st.title("Gerador de Slides da Pauta")

# ================================
# CONFIGURAÇÕES
# ================================

st.subheader("Configurações")

turma = st.text_input(
    "Nome da Turma",
    value="2ª Turma Cível"
)

st.write("### Substituição de nomes dos desembargadores")

col1, col2 = st.columns(2)

with col1:
    nome1 = st.text_input("JOAO EGMONT LEONCIO LOPES","JOÃO EGMONT")
    nome2 = st.text_input("HECTOR VALVERDE SANTANNA","HÉCTOR VALVERDE")

with col2:
    nome3 = st.text_input("RENATO RODOVALHO SCUSSEL","RENATO SCUSSEL")
    nome4 = st.text_input("FERNANDO ANTÔNIO TAVERNARD LIMA","FERNANDO TAVERNARD")

substituicoes = {
    "JOAO EGMONT LEONCIO LOPES": nome1,
    "HECTOR VALVERDE SANTANNA": nome2,
    "RENATO RODOVALHO SCUSSEL": nome3,
    "FERNANDO ANTÔNIO TAVERNARD LIMA": nome4
}

st.divider()

# ================================
# EXTRAIR DOCX
# ================================

def extrair_dados_docx(arquivo):

    doc = Document(arquivo)

    linhas = []

    for p in doc.paragraphs:
        texto = p.text.strip()
        if texto:
            linhas.append(texto)

    for tabela in doc.tables:
        for row in tabela.rows:
            for cell in row.cells:
                texto = cell.text.strip()
                if texto:
                    linhas.append(texto)

    dados = []

    numero = None
    processo = None
    desembargador = None

    for texto in linhas:

        if re.fullmatch(r"\d+", texto):
            numero = texto
            continue

        if re.search(r"\d{7}-\d{2}\.\d{4}\.\d\.\d{2}\.\d{4}", texto):
            processo = texto
            continue

        if texto.isupper() and len(texto) > 5:

            desembargador = texto

            if numero and processo:

                dados.append({
                    "numero": numero,
                    "processo": processo,
                    "desembargador": desembargador,
                    "segredo": "NÃO"
                })

                numero = None
                processo = None
                desembargador = None

    return pd.DataFrame(dados)

# ================================
# TEXTO SLIDE
# ================================

def adicionar_texto(slide, texto, x, y, largura, tamanho, cor):

    caixa = slide.shapes.add_textbox(x, y, largura, Inches(1.2))

    tf = caixa.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    run = p.add_run()
    run.text = texto

    run.font.name = "Century Gothic"
    run.font.size = Pt(tamanho)
    run.font.bold = True
    run.font.color.rgb = cor

# ================================
# UPLOAD DOCX
# ================================

arquivo = st.file_uploader("Enviar DOCX da pauta", type=["docx"])

# ================================
# PROCESSAR
# ================================

if arquivo:

    df = extrair_dados_docx(arquivo)

    if df.empty:
        st.error("Nenhum processo foi encontrado no DOCX")
        st.stop()

    st.success("Pauta extraída")

    st.subheader("Editar pauta")

    df_editado = st.data_editor(
        df,
        column_config={
            "segredo": st.column_config.SelectboxColumn(
                "Segredo",
                options=["NÃO","SIM"]
            )
        },
        use_container_width=True
    )

    if st.button("Gerar apresentação"):

        if not os.path.exists("modelo.pptx"):
            st.error("modelo.pptx não encontrado")
            st.stop()

        prs = Presentation("modelo.pptx")

        largura_slide = prs.slide_width
        slide_ref = prs.slides[1]

        img_ref = None

        for shape in slide_ref.shapes:
            if shape.shape_type == 13:
                img_ref = shape
                break

        layout = prs.slide_layouts[6]

        progress = st.progress(0)

        total = len(df_editado)

        for i,row in df_editado.iterrows():

            slide = prs.slides.add_slide(layout)

            if img_ref:

                img_stream = BytesIO(img_ref.image.blob)

                slide.shapes.add_picture(
                    img_stream,
                    img_ref.left,
                    img_ref.top,
                    img_ref.width,
                    img_ref.height
                )

            adicionar_texto(
                slide,
                turma,
                Inches(3.5),
                Inches(1.15),
                Inches(4),
                18,
                RGBColor(255,255,255)
            )

            proc = str(row["processo"]).split(".8")[0]

            texto = f"{proc} ({row['numero']})"

            adicionar_texto(
                slide,
                texto,
                0,
                Inches(2.7),
                largura_slide,
                60,
                RGBColor(0,0,0)
            )

            nome_original = str(row["desembargador"]).upper().strip()
            nome = substituicoes.get(nome_original, nome_original)

            relator = f"RELATOR:\nDESEMBARGADOR {nome}"

            adicionar_texto(
                slide,
                relator,
                0,
                Inches(4.2),
                largura_slide,
                32,
                RGBColor(0,0,0)
            )

            if row["segredo"] == "SIM":

                adicionar_texto(
                    slide,
                    "SEGREDO DE JUSTIÇA",
                    0,
                    Inches(5.3),
                    largura_slide,
                    40,
                    RGBColor(255,0,0)
                )

            progress.progress((i+1)/total)

        slides_gerados = len(df_editado)

        while len(prs.slides) > slides_gerados + 1:
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[1])

        output = BytesIO()
        prs.save(output)

        st.success("Slides gerados!")

        st.download_button(
            "Baixar PowerPoint",
            output.getvalue(),
            file_name="slides_pauta.pptx"
        )
